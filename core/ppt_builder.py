"""
PowerPoint 빌더 모듈

엑셀 데이터와 템플릿을 결합하여 PowerPoint를 생성하는 메인 로직을 제공합니다.
"""
import pandas as pd
from pptx import Presentation

from utils.text_utils import normalize_key
from utils.slide_utils import duplicate_slide
from core.template_handler import (
    fill_template_with_row, 
    restore_template_shapes, 
    backup_template_shapes
)


def build_ppt_from_excel(template_path, excel_path, out_path, sheet_name=0):
    """
    엑셀 데이터를 기반으로 템플릿에서 PowerPoint 프레젠테이션을 생성합니다.
    
    전체 작동 흐름:
        1. 템플릿 PPT 및 엑셀 파일 로드
        2. 템플릿의 첫 번째 슬라이드를 템플릿으로 사용
        3. 템플릿의 원본 상태를 XML 레벨에서 백업
        4. 엑셀의 각 행에 대해:
           a. 템플릿 슬라이드에 데이터 채우기
           b. 채워진 템플릿 슬라이드를 복제하여 추가
           c. 템플릿을 원본 상태로 복원 (다음 행을 위해)
        5. 원본 템플릿 슬라이드 제거
        6. 최종 PPT 저장
    
    Args:
        template_path (str): 템플릿 PPT 파일 경로
        excel_path (str): 엑셀 데이터 파일 경로
        out_path (str): 출력 PPT 파일 경로
        sheet_name (int or str): 엑셀 시트 이름 또는 인덱스 (기본값: 0)
        
    Returns:
        str: 생성된 출력 파일 경로
        
    Raises:
        RuntimeError: 엑셀에 데이터가 없는 경우
        
    템플릿 요구사항:
        - 첫 번째 슬라이드가 템플릿으로 사용됨
        - 템플릿의 텍스트 박스에 엑셀 컬럼명을 정확히 입력
        - 이미지는 컬럼명을 '_IMG'로 끝내거나 셀에 이미지 파일 경로 입력
        
    엑셀 요구사항:
        - 첫 행은 컬럼명 (템플릿의 텍스트와 매칭됨)
        - 각 행은 하나의 슬라이드로 변환됨
        - 빈 셀은 빈 텍스트로 처리됨
        
    Example:
        템플릿 (template.pptx):
            슬라이드 1:
                [고객명] [전화번호] [로고_IMG]
                
        엑셀 (data.xlsx):
            | 고객명  | 전화번호      | 로고_IMG           |
            |--------|--------------|-------------------|
            | 홍길동 | 010-1234-5678| images/logo1.png  |
            | 김철수 | 010-9876-5432| images/logo2.png  |
            
        결과 (output.pptx):
            슬라이드 1: [홍길동] [010-1234-5678] [이미지: logo1.png]
            슬라이드 2: [김철수] [010-9876-5432] [이미지: logo2.png]
            
    성능 최적화:
        - 템플릿을 한 번만 로드하고 메모리에서 재사용
        - XML 레벨 복사로 빠른 슬라이드 복제
        - 관계(Relationship) 자동 리매핑으로 이미지 정확도 보장
        
    Note:
        - 템플릿 슬라이드는 최종 PPT에서 제거됨
        - 그룹 Shape 내부의 텍스트도 자동으로 처리됨
        - 이미지 삽입 실패 시 경고 출력 후 계속 진행
    """
    # 1. 파일 로드
    print(f"템플릿 로드 중: {template_path}")
    print(f"엑셀 로드 중: {excel_path}")
    
    prs = Presentation(template_path)
    df = pd.read_excel(excel_path, sheet_name=sheet_name)
    
    # 2. 데이터 유효성 검증
    if df.empty:
        raise RuntimeError('엑셀에 데이터가 없습니다.')

    # 3. 템플릿 슬라이드 준비 (첫 번째 슬라이드)
    tmpl = prs.slides[0]

    # 디버그 정보 출력
    print(f"\n템플릿 슬라이드 shape 수: {len(tmpl.shapes)}")
    print(f"엑셀 열: {list(df.columns)}")
    print(f"엑셀 행 수: {len(df)}")
    print("\n템플릿 텍스트 박스들:")
    for sh in tmpl.shapes:
        if getattr(sh, 'has_text_frame', False):
            text = normalize_key(sh.text)
            print(f"  - '{text}'")

    # 4. 템플릿의 원본 상태 백업
    print("\n템플릿 백업 중...")
    original_shapes_xml = backup_template_shapes(tmpl)

    # 5. 각 행마다 슬라이드 생성
    print(f"\n총 {len(df)}개의 슬라이드 생성 시작...\n")
    
    for idx, row in df.iterrows():
        print(f"슬라이드 {idx + 1}/{len(df)} 생성 중...")
        
        # a. 템플릿 슬라이드에 데이터 채우기
        fill_template_with_row(tmpl, row)
        
        # b. 채워진 템플릿을 복제하여 새 슬라이드 생성
        duplicate_slide(prs, tmpl)
        
        # c. 템플릿을 원본 상태로 복원 (다음 행을 위해)
        restore_template_shapes(tmpl, original_shapes_xml)

    # 6. 원본 템플릿 슬라이드 제거
    print("\n원본 템플릿 슬라이드 제거 중...")
    try:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    except Exception as e:
        print(f"템플릿 슬라이드 제거 실패 (무시 가능): {e}")

    # 7. 최종 파일 저장
    print(f"\n파일 저장 중: {out_path}")
    prs.save(out_path)
    
    print(f"✓ 완료: {len(df)}개의 슬라이드가 생성되었습니다.")
    return out_path


def validate_template(template_path):
    """
    템플릿 파일의 유효성을 검증하고 정보를 출력합니다.
    
    Args:
        template_path (str): 템플릿 PPT 파일 경로
        
    Returns:
        dict: 템플릿 정보 딕셔너리
            - slide_count: 슬라이드 수
            - text_boxes: 첫 번째 슬라이드의 텍스트 박스 목록
            
    Example:
        >>> info = validate_template("template.pptx")
        >>> print(f"텍스트 박스: {info['text_boxes']}")
    """
    prs = Presentation(template_path)
    
    if len(prs.slides) == 0:
        raise RuntimeError("템플릿에 슬라이드가 없습니다.")
    
    tmpl = prs.slides[0]
    text_boxes = []
    
    for shape in tmpl.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = normalize_key(shape.text)
            if text:
                text_boxes.append(text)
    
    info = {
        'slide_count': len(prs.slides),
        'text_boxes': text_boxes
    }
    
    print(f"템플릿 슬라이드 수: {info['slide_count']}")
    print(f"텍스트 박스 ({len(text_boxes)}개):")
    for tb in text_boxes:
        print(f"  - {tb}")
    
    return info
