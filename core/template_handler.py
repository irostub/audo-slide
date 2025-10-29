"""
템플릿 슬라이드 처리 모듈

템플릿 슬라이드에 데이터를 채우고 복원하는 기능을 제공합니다.
"""
import copy
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE_TYPE

from utils.text_utils import normalize_key, get_text_content, set_text_preserve_style
from utils.image_utils import fill_picture, IMAGE_EXTENSIONS


def fill_template_with_row(tmpl_slide, row, image_suffixes=IMAGE_EXTENSIONS):
    """
    템플릿 슬라이드에 엑셀 행 데이터를 채웁니다.
    
    작동 방식:
        1. 엑셀 행의 모든 컬럼을 정규화된 키-값 딕셔너리로 변환
        2. 템플릿의 모든 Shape를 순회 (그룹 내부 Shape 포함)
        3. 각 Shape의 텍스트를 키로 사용하여 엑셀 데이터와 매칭
        4. 매칭된 데이터를 Shape에 채움 (텍스트 또는 이미지)
    
    Args:
        tmpl_slide: 템플릿 슬라이드 객체
        row: pandas Series (엑셀의 한 행)
        image_suffixes: 이미지 파일로 인식할 확장자 튜플
        
    매칭 규칙:
        - 엑셀 컬럼명과 Shape의 텍스트를 정규화하여 비교
        - 정확히 일치하는 경우만 데이터 채움
        
    데이터 타입별 처리:
        1. None 또는 NaN: 빈 문자열로 채움
        2. 이미지: 컬럼명이 '_IMG'로 끝나거나 값이 이미지 확장자인 경우
           - 해당 경로의 이미지를 Shape에 삽입
        3. 일반 텍스트: 문자열로 변환하여 Shape에 채움 (스타일 유지)
        
    그룹 Shape 처리:
        - 그룹 내부의 모든 Shape도 개별적으로 처리
        - 그룹 구조는 유지되며 내부 데이터만 변경
        
    Example:
        엑셀:
            | 고객명    | 금액      | 로고_IMG           |
            |----------|----------|-------------------|
            | 홍길동   | 1000000  | images/logo.png   |
        
        템플릿:
            [고객명] [금액] [로고_IMG]
            
        결과:
            [홍길동] [1000000] [이미지: logo.png]
            
    Note:
        - 이 함수는 템플릿을 직접 수정합니다 (in-place)
        - 복제 전에 호출하여 데이터를 채운 후 복제해야 합니다
        - 매칭되지 않는 Shape는 변경되지 않습니다
    """
    # 1. 엑셀 행 데이터를 정규화된 딕셔너리로 변환
    row_dict = {}
    for k in row.index:
        k_norm = normalize_key(str(k))
        if k_norm:  # 빈 키는 무시
            row_dict[k_norm] = row[k]

    print(f"  엑셀 원본 컬럼: {list(row.index)}")
    print(f"  엑셀 정규화 키: {list(row_dict.keys())}")

    # 2. 템플릿의 모든 Shape 수집 (그룹 내부 Shape 포함)
    shapes_to_process = []
    try:
        for shape in tmpl_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 그룹인 경우 내부 Shape들을 펼침
                try:
                    shapes_to_process.extend(shape.shapes)
                except Exception:
                    # 그룹 처리 실패 시 무시
                    pass
            else:
                shapes_to_process.append(shape)
    except Exception as e:
        print(f"  경고: Shape 수집 중 오류: {e}")
        return

    # 3. 매칭 통계
    matched_count = 0
    unmatched_shapes = []

    # 4. 각 Shape 처리
    for shape in shapes_to_process:
        try:
            # Shape의 텍스트를 키로 사용
            key_raw = get_text_content(shape)
            key = normalize_key(key_raw)

            # 디버깅: 원본과 정규화된 텍스트 비교
            if key_raw and key_raw != key:
                print(f"    [정규화] '{key_raw}' -> '{key}'")

            # 엑셀 데이터와 매칭되지 않으면 스킵
            if not key:
                continue

            if key not in row_dict:
                unmatched_shapes.append(key)
                continue

            val = row_dict[key]

            # None 또는 빈 값 처리
            # dtype=str로 읽었기 때문에 'nan' 문자열로 올 수 있음
            if val is None or (isinstance(val, str) and val.lower() == 'nan'):
                print(f"    매칭: '{key}' -> (빈 값)")
                if getattr(shape, 'has_text_frame', False):
                    set_text_preserve_style(shape, '')
                matched_count += 1
                continue

            # 이미 문자열로 읽혔으므로 str() 변환 필요 없음
            sval = str(val).strip()

            # 빈 문자열 처리
            if not sval:
                print(f"    매칭: '{key}' -> (빈 문자열)")
                if getattr(shape, 'has_text_frame', False):
                    set_text_preserve_style(shape, '')
                matched_count += 1
                continue

            print(f"    매칭: '{key}' -> '{sval}'")

            # 이미지 처리: 키가 '_IMG'로 끝나거나 값이 이미지 확장자인 경우
            if key.upper().endswith('_IMG') or sval.lower().endswith(image_suffixes):
                fill_picture(shape, sval)
                matched_count += 1
                continue

            # 일반 텍스트 처리 (스타일 유지)
            set_text_preserve_style(shape, sval)
            matched_count += 1

        except Exception as e:
            print(f"  경고: Shape 처리 중 오류: {e}")
            continue

    # 5. 매칭 결과 요약
    if unmatched_shapes:
        print(f"  ℹ 매칭 안됨 ({len(unmatched_shapes)}개): {unmatched_shapes[:5]}")
        if len(unmatched_shapes) > 5:
            print(f"    ... 외 {len(unmatched_shapes) - 5}개")
    print(f"  ✓ 총 {matched_count}개 필드 채움 완료")


def restore_template_shapes(tmpl_slide, original_shapes_xml):
    """
    템플릿 슬라이드의 Shape들을 원본 XML로 복원합니다.
    
    이 함수는 데이터를 채운 템플릿을 다시 깨끗한 상태로 되돌립니다.
    이를 통해 하나의 템플릿으로 여러 슬라이드를 생성할 수 있습니다.
    
    작동 방식:
        1. 현재 슬라이드의 모든 Shape 제거
        2. 백업해둔 원본 Shape XML을 다시 삽입
        
    Args:
        tmpl_slide: 템플릿 슬라이드 객체
        original_shapes_xml: 원본 Shape XML 요소 리스트 (deep copy된 것)
        
    사용 패턴:
        # 1. 원본 백업
        original_xml = [copy.deepcopy(shape.element) for shape in tmpl.shapes]
        
        # 2. 데이터 채우기
        fill_template_with_row(tmpl, row1)
        
        # 3. 복제
        slide1 = duplicate_slide(prs, tmpl)
        
        # 4. 템플릿 복원 (다음 행을 위해)
        restore_template_shapes(tmpl, original_xml)
        
        # 5. 반복...
        
    Why?
        - 매번 템플릿 파일을 다시 로드하는 것보다 훨씬 빠름
        - 메모리에서 XML만 교체하므로 효율적
        - 원본 템플릿의 모든 속성(스타일, 애니메이션 등)을 완벽히 보존
        
    Note:
        - original_shapes_xml은 반드시 deep copy여야 함
        - 얕은 복사를 사용하면 원본이 손상될 수 있음
    """
    # 1. 기존 Shape 모두 제거
    for shape in list(tmpl_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)
    
    # 2. 원본 Shape XML을 다시 삽입
    for orig_xml in original_shapes_xml:
        new_el = copy.deepcopy(orig_xml)
        tmpl_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def backup_template_shapes(tmpl_slide):
    """
    템플릿 슬라이드의 모든 Shape를 XML 레벨에서 백업합니다.
    
    Args:
        tmpl_slide: 템플릿 슬라이드 객체
        
    Returns:
        list: 모든 Shape의 deep copy된 XML 요소 리스트
        
    Example:
        >>> backup = backup_template_shapes(template_slide)
        >>> # ... 템플릿 수정 작업 ...
        >>> restore_template_shapes(template_slide, backup)
    """
    return [copy.deepcopy(shape.element) for shape in tmpl_slide.shapes]
